#!/usr/bin/env python3
"""
extract_content.py - 从 PDF / DOCX / PPTX 文件中提取文本内容。

统一入口，根据文件扩展名自动选择提取方式：
  - PDF:  使用 pdfplumber
  - DOCX: 优先使用 pandoc，回退到 unpack.py + XML 文本提取
  - PPTX: 优先使用 markitdown，回退到 python-pptx

Usage:
    python3 extract_content.py <file_path> [--output <output_file>] [--max-pages 50]

Dependencies (按需安装):
    pip install pdfplumber        # PDF 解析
    pip install markitdown        # PPTX 解析（首选）
    pip install python-pptx       # PPTX 解析（备选）
    系统安装 pandoc               # DOCX 解析（首选）
"""

from __future__ import annotations

import argparse
import os
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree


# ---------------------------------------------------------------------------
# PDF 提取
# ---------------------------------------------------------------------------
def extract_pdf(file_path: str, max_pages: int) -> str:
    """使用 pdfplumber 提取 PDF 文本内容。"""
    try:
        import pdfplumber  # noqa: F811 # type: ignore
    except ImportError:
        print("Error: pdfplumber 未安装，请执行: pip install pdfplumber", file=sys.stderr)
        sys.exit(1)

    all_text: list[str] = []
    empty_page_count = 0

    with pdfplumber.open(file_path) as pdf:
        total = min(len(pdf.pages), max_pages)
        if len(pdf.pages) > max_pages:
            print(f"提示: 文档共 {len(pdf.pages)} 页，仅提取前 {max_pages} 页（可通过 --max-pages 调整）",
                  file=sys.stderr)

        for i, page in enumerate(pdf.pages[:total], 1):
            text = page.extract_text()
            if text and text.strip():
                all_text.append(f"--- Page {i} ---\n{text}")
            else:
                all_text.append(f"--- Page {i} ---\n[本页无可提取文本]")
                empty_page_count += 1

    # 扫描件检测：超过一半页面无文本则提示
    if total > 0 and empty_page_count > total * 0.5:
        print(
            f"警告: {empty_page_count}/{total} 页无可提取文本，该 PDF 可能是扫描件。"
            "建议使用 Claude 视觉能力直接读取图片，或先进行 OCR 处理。",
            file=sys.stderr,
        )

    return "\n\n".join(all_text)


# ---------------------------------------------------------------------------
# DOCX 提取
# ---------------------------------------------------------------------------
def _extract_docx_pandoc(file_path: str) -> str | None:
    """尝试使用 pandoc 将 DOCX 转为纯文本。"""
    if not shutil.which("pandoc"):
        return None
    try:
        result = subprocess.run(
            ["pandoc", "-t", "plain", "--wrap=none", file_path],
            capture_output=True, text=True, timeout=60, check=True,
        )
        return result.stdout
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
        return None


def _extract_docx_xml(file_path: str) -> str:
    """回退方案：解压 DOCX 后从 document.xml 提取纯文本。"""
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: list[str] = []

    with zipfile.ZipFile(file_path, "r") as z:
        if "word/document.xml" not in z.namelist():
            print("Error: DOCX 文件结构异常，未找到 word/document.xml", file=sys.stderr)
            sys.exit(1)
        with z.open("word/document.xml") as f:
            tree = ElementTree.parse(f)  # noqa: S314

    for para in tree.iter(f"{{{ns['w']}}}p"):
        texts = [t.text for t in para.iter(f"{{{ns['w']}}}t") if t.text]
        if texts:
            paragraphs.append("".join(texts))

    return "\n\n".join(paragraphs)


def extract_docx(file_path: str) -> str:
    """提取 DOCX 文本内容，优先使用 pandoc。"""
    text = _extract_docx_pandoc(file_path)
    if text and text.strip():
        return text

    if not shutil.which("pandoc"):
        print("提示: 未检测到 pandoc，使用内置 XML 解析（格式保留较少）。"
              "推荐安装 pandoc 以获得更好效果: https://pandoc.org/installing.html",
              file=sys.stderr)

    return _extract_docx_xml(file_path)


# ---------------------------------------------------------------------------
# PPTX 提取
# ---------------------------------------------------------------------------
def _extract_pptx_markitdown(file_path: str) -> str | None:
    """尝试使用 markitdown 提取 PPTX 内容。"""
    try:
        from markitdown import MarkItDown  # type: ignore[import-untyped]
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content
    except ImportError:
        return None
    except Exception:
        return None


def _extract_pptx_python_pptx(file_path: str) -> str | None:
    """回退方案：使用 python-pptx 提取 PPTX 文本。"""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return None

    prs = Presentation(file_path)
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
        slides_text.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    return "\n\n".join(slides_text)


def extract_pptx(file_path: str) -> str:
    """提取 PPTX 文本内容，优先使用 markitdown。"""
    text = _extract_pptx_markitdown(file_path)
    if text and text.strip():
        return text

    text = _extract_pptx_python_pptx(file_path)
    if text and text.strip():
        return text

    print(
        "Error: 无法解析 PPTX 文件，请安装以下任一依赖:\n"
        "  pip install markitdown    （推荐）\n"
        "  pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------
SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf": "PDF",
    ".docx": "DOCX",
    ".pptx": "PPTX",
}

IMAGE_EXTENSIONS: set[str] = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg"}


def main() -> None:
    """命令行入口：根据文件扩展名分发到对应的提取函数。"""
    parser = argparse.ArgumentParser(
        description="从 PDF / DOCX / PPTX 文件中提取文本内容",
    )
    parser.add_argument("file_path", help="待提取的文件路径")
    parser.add_argument("--output", "-o", help="输出文件路径（默认输出到 stdout）")
    parser.add_argument("--max-pages", type=int, default=50,
                        help="PDF 最大提取页数（默认 50）")
    args = parser.parse_args()

    file_path = args.file_path
    if not os.path.isfile(file_path):
        print(f"Error: 文件不存在: {file_path}", file=sys.stderr)
        sys.exit(1)

    ext = Path(file_path).suffix.lower()

    # 图片文件提示使用 Claude 视觉能力
    if ext in IMAGE_EXTENSIONS:
        print(
            f"提示: {ext} 是图片文件，无需通过此脚本提取。\n"
            "请使用 Claude 的 Read 工具直接读取图片，Claude 视觉能力可识别技术架构图、流程图等内容。",
            file=sys.stderr,
        )
        sys.exit(0)

    if ext not in SUPPORTED_EXTENSIONS:
        print(
            f"Error: 不支持的文件类型 '{ext}'。\n"
            f"支持的类型: {', '.join(SUPPORTED_EXTENSIONS.keys())}",
            file=sys.stderr,
        )
        sys.exit(1)

    # 根据扩展名分发提取
    extractors = {
        ".pdf": lambda: extract_pdf(file_path, args.max_pages),
        ".docx": lambda: extract_docx(file_path),
        ".pptx": lambda: extract_pptx(file_path),
    }

    try:
        text = extractors[ext]()
    except Exception as e:
        print(f"Error: 处理 {SUPPORTED_EXTENSIONS[ext]} 文件时出错: {e}", file=sys.stderr)
        sys.exit(1)

    # 输出结果
    if args.output:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(text, encoding="utf-8")
        print(f"内容已提取到: {args.output}")
    else:
        print(text)


if __name__ == "__main__":
    main()
